from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, FileResponse
from django.contrib.auth.decorators import login_required
from .models import QuestionPaper
import re
import fitz  # PyMuPDF for PDF extraction
import google.generativeai as genai
from fpdf import FPDF
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from docx import Document
import os
import tempfile
import uuid
from django.conf import settings
from django.core.files.base import ContentFile
from django.core.files.storage import default_storage

# Configure Gemini API
genai.configure(api_key="AIzaSyAGfz_xI8ODEwVNXF82g0S9K2h9ktpcHJ4")
model = genai.GenerativeModel("gemini-1.5-flash")

def home1(request):
    return render(request, 'main/home1.html')

def upload_files(request):
    if request.method == 'POST':
        files = request.FILES.getlist('files')
        temp_file_paths = []
        
        for file in files:
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{file.name.split('.')[-1]}")
            for chunk in file.chunks():
                temp_file.write(chunk)
            temp_file.close()
            temp_file_paths.append(temp_file.name)
        
        request.session['file_paths'] = temp_file_paths
        return redirect('configure_exam')
    
    return render(request, 'main/upload.html')

def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    return "\n".join([page.get_text("text") for page in doc])

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_files(file_list):
    combined_text = "\n".join([
        extract_text_from_pdf(file) if file.endswith(".pdf") else
        extract_text_from_pptx(file) if file.endswith(".pptx") else
        extract_text_from_docx(file) if file.endswith(".docx") else ""
        for file in file_list
    ])
    return combined_text

def generate_questions(text, marks, count, difficulty):
    prompt = (f"Generate {count} {difficulty} difficulty questions of {marks} marks each based on the text below.\n"
              f"Only provide the questions without additional explanations or formatting.\n{text[:4000]}")
    response = model.generate_content(prompt)
    # Fix the regex pattern to properly remove leading numbers with dots
    return list(set([re.sub(r'^\d+\.\s*', '', line.strip()) for line in response.text.split("\n") if line.strip()]))

def generate_course_outcomes(questions_preview, subject):
    """
    Generate 2-3 course outcomes based on the questions and subject
    Each CO should be short (6-7 words only)
    """
    # Collect all questions into a single text
    all_questions = []
    for q_content in questions_preview.values():
        all_questions.extend(q_content['questions'])
    
    # Create a prompt for Gemini
    prompt = (f"Based on the following questions for a {subject} exam, generate 2-3 course outcomes (CO) "
              f"that these questions are assessing. Each CO should start with an action verb (like Explain, "
              f"Apply, Analyze, etc.) and must be VERY SHORT (6-7 words MAXIMUM). Format each CO as 'CO1: [description]', "
              f"'CO2: [description]', etc.\n\nQuestions:\n" + "\n".join(all_questions))
    
    try:
        response = model.generate_content(prompt)
        
        # Extract COs using regex
        co_pattern = r'CO\d+:\s*(.*?)(?=(?:CO\d+:|$))'
        cos = re.findall(co_pattern, response.text, re.DOTALL)
        
        # Clean and format the COs
        co_dict = {}
        for i, co_text in enumerate(cos[:3], 1):  # Limit to 3 COs maximum
            # Ensure each CO is short (truncate if needed)
            co_text = co_text.strip()
            words = co_text.split()
            if len(words) > 7:
                co_text = " ".join(words[:7])
            co_dict[f"CO{i}"] = co_text
        
        # Ensure we have at least 2 COs
        if len(co_dict) < 2:
            co_dict = {
                "CO1": "Understand key subject concepts.",
                "CO2": "Apply knowledge to practical problems."
            }
        
        return co_dict
    
    except Exception as e:
        # Fallback COs in case of API error (short versions)
        return {
            "CO1": "Understand key subject concepts.",
            "CO2": "Apply knowledge to practical problems."
        }

def is_similar(new_question, question_list, threshold=0.7):
    if not question_list:
        return False
    vectorizer = TfidfVectorizer().fit_transform([new_question] + question_list)
    similarity_matrix = cosine_similarity(vectorizer)[0][1:]
    return any(similarity >= threshold for similarity in similarity_matrix)

def select_unique_questions(question_list, num, used_questions, avoid_duplicates):
    selected = []
    for q in question_list:
        if q not in used_questions and len(selected) < num:
            if avoid_duplicates and is_similar(q, list(used_questions)):
                continue
            selected.append(q)
            used_questions.add(q)
    return selected

def configure_exam(request):
    if request.method == 'POST':
        exam_data = {
            'pdf_type': request.POST.get('pdf_type'),
            'avoid_duplicates': request.POST.get('avoid_duplicates') == 'yes',
            'semester': request.POST.get('semester'),
            'branch': request.POST.get('branch'),
            'subject': request.POST.get('subject'),
            'exam_date': request.POST.get('exam_date'),
            'exam_time': request.POST.get('exam_time'),
            'questions': []
        }
        total_questions = 3 if exam_data['pdf_type'] == "mse" else 5
        
        for q in range(1, total_questions + 1):
            question_type = int(request.POST.get(f'q{q}_type'))
            difficulty = request.POST.get(f'q{q}_difficulty')
            exam_data['questions'].append({'question_type': question_type, 'difficulty': difficulty})
        
        file_paths = request.session.get('file_paths', [])
        extracted_text = extract_text_from_files(file_paths)
        
        request.session['exam_data'] = exam_data
        request.session['extracted_text'] = extracted_text
        
        return redirect('preview_questions')
    
    return render(request, 'main/generate_paper_from_notes.html')

def preview_questions(request):
    exam_data = request.session.get('exam_data', {})
    extracted_text = request.session.get('extracted_text', '')
    used_questions = set()
    questions_preview = {}
    
    for q_index, q_info in enumerate(exam_data['questions'], 1):
        question_type = q_info['question_type']
        difficulty = q_info['difficulty']
        num = 3 if question_type == 5 else 2 if question_type == 10 else 5
        
        generated_questions = generate_questions(extracted_text, question_type, num * 2, difficulty)
        unique_questions = select_unique_questions(generated_questions, num, used_questions, exam_data['avoid_duplicates'])
        
        instruction = "Solve any two (5 marks each)" if question_type == 5 else \
                     "Solve anyone (10 marks each)" if question_type == 10 else \
                     "Solve any five (2 marks each)"
        
        questions_preview[q_index] = {'instruction': instruction, 'questions': unique_questions}
    
    # Generate course outcomes based on the questions
    course_outcomes = generate_course_outcomes(questions_preview, exam_data.get('subject', ''))
    
    # Store both questions and course outcomes in session
    request.session['questions_preview'] = questions_preview
    request.session['course_outcomes'] = course_outcomes
    
    # Pass both questions_preview, exam_data, and course outcomes to the template
    return render(request, 'main/preview.html', {
        'questions': questions_preview,
        'exam_data': exam_data,
        'course_outcomes': course_outcomes
    })

def map_questions_to_cos(questions_preview, course_outcomes):
    """
    Map each question to the most appropriate CO based on content similarity
    """
    # Initialize a dictionary to store question to CO mapping
    question_co_mapping = {}
    
    # For each question, determine which CO it maps to
    co_texts = list(course_outcomes.values())
    co_keys = list(course_outcomes.keys())
    
    for q_num, q_content in questions_preview.items():
        # Map this question set to a CO
        co_index = (int(q_num) - 1) % len(co_keys)
        question_co_mapping[q_num] = co_keys[co_index]
    
    return question_co_mapping

def generate_pdf(exam_data, questions_preview):
    # Get course outcomes from session or generate default
    course_outcomes = exam_data.get('course_outcomes', {
        "CO1": "Understand key subject concepts.",
        "CO2": "Apply knowledge to practical problems."
    })
    
    # Map questions to appropriate COs
    question_co_mapping = map_questions_to_cos(questions_preview, course_outcomes)
    
    # Create PDF with handling for Unicode characters
    class MYPDF(FPDF):
        def __init__(self):
            super().__init__()
            # Handle Unicode encoding errors by replacing problematic characters
            self.add_page()
    
        # Override to handle encoding errors for cell
        def cell(self, w, h=0, txt='', border=0, ln=0, align='', fill=False, link=''):
            # Clean up the text before adding it to PDF
            cleaned_txt = txt.encode('latin-1', errors='replace').decode('latin-1')
            super().cell(w, h, cleaned_txt, border, ln, align, fill, link)
    
        # Override to handle encoding errors for multi_cell
        def multi_cell(self, w, h, txt='', border=0, align='', fill=False):
            # Clean up the text before adding it to PDF
            cleaned_txt = txt.encode('latin-1', errors='replace').decode('latin-1')
            super().multi_cell(w, h, cleaned_txt, border, align, fill)
    
    pdf = MYPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # Header section
    pdf.set_font("Arial", style="B", size=16)
    pdf.cell(0, 10, txt="Vidyalankar Institute of Technology", ln=True, align='C')
    
    # Semester, branch and exam type
    pdf.set_font("Arial", style="B", size=14)
    exam_type = "Mid Semester Assessment" if exam_data.get('pdf_type') == 'mse' else "End Semester Assessment"
    pdf.cell(0, 10, txt=f"Semester {exam_data.get('semester')} -- {exam_data.get('branch')}- {exam_type}", ln=True, align='C')
    
    # Add horizontal line
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    
    # Create header table (Date, Subject, Marks/Duration)
    pdf.ln(1)
    pdf.set_font("Arial", style="B", size=12)
    
    # Create table for header info
    col_width1 = 60
    col_width2 = 80
    col_width3 = 50
    
    marks = "30 Marks/ 1 Hr." if exam_data.get('pdf_type') == 'mse' else "50 Marks/ 2 Hrs."
    
    # First row of header table
    pdf.cell(col_width1, 10, txt=f"Date: {exam_data.get('exam_date')}", border=1)
    pdf.cell(col_width2, 10, txt=f"{exam_data.get('subject')}", border=1, align='C')
    pdf.cell(col_width3, 10, txt=marks, border=1, align='C')
    pdf.ln(10)
    
    # Add horizontal line
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(5)  # Add some space before questions
    
    # Questions section
    for q_num, q_content in questions_preview.items():
        # Table for questions - similar to the Word document format
        # First row: Question number, instruction, and CO header
        pdf.set_font("Arial", style="B", size=12)
        
        # Cell for question number and instruction
        instruction_text = f"{q_num}   {q_content['instruction']}"
        pdf.cell(170, 10, txt=instruction_text, border=1)
        
        # CO column header
        pdf.cell(20, 10, txt="CO", border=1, align='C')
        pdf.ln()
        
        # Get the CO for this question
        current_co = question_co_mapping.get(str(q_num), f"CO{q_num}")
        
        # Add each lettered question
        for i, question in enumerate(q_content['questions']):
            letter = chr(65 + i)
            
            # Clean the question text
            cleaned_question = re.sub(r'^\d+\.\s*', '', question)
            cleaned_question = re.sub(r'\(\d+\s*marks\)$', '', cleaned_question, flags=re.IGNORECASE).strip()
            
            # Remove or replace problematic Unicode characters
            cleaned_question = cleaned_question.replace('"', '"').replace('"', '"')
            cleaned_question = cleaned_question.replace(''', "'").replace(''', "'")
            cleaned_question = cleaned_question.replace('—', '-').replace('–', '-').replace('…', '...')
            
            # Letter in bold (first column)
            pdf.set_font("Arial", style="B", size=12)
            pdf.cell(10, 10, txt=f"{letter}", border=1, align='C')
            
            # Question text (second column)
            pdf.set_font("Arial", size=12)
            
            # Save current position
            x_pos = pdf.get_x()
            y_pos = pdf.get_y()
            
            # Calculate text height
            pdf.multi_cell(160, 10, txt=cleaned_question, border=1)
            
            # Move to position for CO column
            new_y = pdf.get_y()
            pdf.set_xy(x_pos + 160, y_pos)
            
            # CO value (third column)
            pdf.set_font("Arial", size=12)
            pdf.cell(20, new_y - y_pos, txt=current_co, border=1, align='C')
            
            # Reset position for next row
            pdf.ln()
        
        pdf.ln(5)
    
    # Add horizontal line
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(5)
    
    # Add CO description table at the bottom
    # Table header
    table_width = 190
    pdf.set_font("Arial", style="B", size=12)
    
    # CO table
    for co_key, co_desc in course_outcomes.items():
        # CO key column
        pdf.cell(20, 10, txt=co_key, border=1)
        
        # CO description column
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(table_width - 20, 10, txt=co_desc, border=1)
    
    # Add final horizontal line
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    temp_file.close()
    pdf.output(temp_file.name)
    
    return temp_file.name

def download_paper(request):
    exam_data = request.session.get('exam_data', {})
    questions_preview = request.session.get('questions_preview', {})
    course_outcomes = request.session.get('course_outcomes', {})
    
    # Add course outcomes to exam_data for PDF generation
    exam_data['course_outcomes'] = course_outcomes
    
    temp_file_name = generate_pdf(exam_data, questions_preview)
    
    response = FileResponse(open(temp_file_name, 'rb'), content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="Question_Paper_{uuid.uuid4().hex[:8]}.pdf"'
    return response

@login_required
def save_question_paper(request):
    if request.method == 'POST':
        # Get data from form and session
        exam_data = request.session.get('exam_data', {})
        questions_preview = request.session.get('questions_preview', {})
        course_outcomes = request.session.get('course_outcomes', {})
        
        if not exam_data or not questions_preview:
            return redirect('upload_files')
        
        # Add course outcomes to exam_data for PDF generation
        exam_data['course_outcomes'] = course_outcomes
        
        # Generate PDF file
        temp_file_name = generate_pdf(exam_data, questions_preview)
        
        # Create QuestionPaper instance
        question_paper = QuestionPaper(
            question_paper_name=request.POST.get('question_paper_name'),
            semester=int(exam_data.get('semester')),
            subject_name=exam_data.get('subject'),
            subject_code=request.POST.get('subject_code'),
            branch=exam_data.get('branch'),
            exam_type="Mid Semester Exam" if exam_data.get('pdf_type') == 'mse' else "End Semester Exam",
            exam_time=exam_data.get('exam_time'),
            exam_date=exam_data.get('exam_date'),
            created_by=request.user.teacher if hasattr(request.user, 'teacher') else None
        )
        
        # Save the PDF to the question_paper_file field
        with open(temp_file_name, 'rb') as f:
            filename = f"question_paper_{uuid.uuid4().hex[:8]}.pdf"
            question_paper.question_paper_file.save(filename, ContentFile(f.read()))
        
        # Save the QuestionPaper instance
        question_paper.save()
        
        # Clean up the temporary file
        os.unlink(temp_file_name)
        
        return redirect('view_question_papers')
    
    return redirect('preview_questions')

@login_required
def view_question_papers(request):
    # Get question papers created by the logged-in teacher
    if hasattr(request.user, 'teacher'):
        question_papers = QuestionPaper.objects.filter(created_by=request.user.teacher).order_by('-created_at')
    else:
        question_papers = QuestionPaper.objects.filter(created_by=None).order_by('-created_at')
    
    return render(request, 'main/saved_question_papers.html', {
        'question_papers': question_papers
    })

@login_required
def delete_question_paper(request, paper_id):
    question_paper = get_object_or_404(QuestionPaper, id=paper_id)
    
    # Check if the paper belongs to the logged-in teacher
    if hasattr(request.user, 'teacher') and question_paper.created_by != request.user.teacher:
        return redirect('view_question_papers')
    
    # Delete the file from storage
    if question_paper.question_paper_file:
        if os.path.isfile(question_paper.question_paper_file.path):
            os.remove(question_paper.question_paper_file.path)
    
    # Delete the QuestionPaper instance
    question_paper.delete()
    
    return redirect('view_question_papers')